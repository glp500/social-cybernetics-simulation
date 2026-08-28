#!/usr/bin/env python3
"""
Simple script to convert the markdown presentation to PowerPoint.

Requires: python-pptx
Install: pip install python-pptx

Usage: python convert_to_pptx.py model-architecture-presentation.md output.pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def create_presentation_from_markdown(md_file, output_pptx):
    """Convert markdown presentation to PowerPoint."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Read markdown file
    with open(md_file, 'r') as f:
        content = f.read()
    
    # Split into slides (separated by "---")
    slides_content = content.split('---')
    
    # Process each slide
    for slide_content in slides_content:
        if not slide_content.strip():
            continue
            
        # Create slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        # Split into lines
        lines = slide_content.strip().split('\n')
        
        # Find title (first line starting with #)
        title = ""
        body_lines = []
        in_title = False
        
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
                in_title = True
            elif line.startswith('## '):
                # Subtitle
                if body_lines and not body_lines[-1].endswith('\n'):
                    body_lines.append('\n')
                body_lines.append(line[3:].strip() + '\n')
            elif line.startswith('### '):
                # Section header
                if body_lines and not body_lines[-1].endswith('\n\n'):
                    body_lines.append('\n')
                body_lines.append(line[4:].strip() + '\n\n')
            elif line.strip():
                body_lines.append(line + '\n')
        
        # Set title
        if title:
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            
            # Format title
            title_text_frame = title_placeholder.text_frame
            title_text_frame.paragraphs[0].font.bold = True
            title_text_frame.paragraphs[0].font.size = Pt(44)
        
        # Set content
        if body_lines:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = ''.join(body_lines)
            
            # Format content
            text_frame = content_placeholder.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.name = 'Calibri'
    
    # Save presentation
    prs.save(output_pptx)
    print(f"✅ Presentation saved to: {output_pptx}")
    print(f"   File size: {len(content)} characters")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_to_pptx.py input.md output.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    create_presentation_from_markdown(input_file, output_file)
